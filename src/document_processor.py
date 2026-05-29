"""文档处理模块 - 处理PDF/DOC/XLS/PPT等文档的内容提取和缩略图生成"""

import io
import csv
from PIL import Image, ImageDraw

# 尝试导入可选依赖
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import xlrd
    HAS_XLRD = True
except ImportError:
    HAS_XLRD = False


class DocumentProcessor:
    """文档处理器 - 提取文档内容并生成缩略图"""
    
    @staticmethod
    def get_font(size):
        """获取字体"""
        try:
            from PIL import ImageFont
            return ImageFont.truetype('simhei.ttf', size)
        except:
            try:
                from PIL import ImageFont
                return ImageFont.truetype('msyh.ttc', size)
            except:
                from PIL import ImageFont
                return ImageFont.load_default()
    
    @staticmethod
    def extract_pdf_content(content):
        """提取PDF内容"""
        if not HAS_PYMUPDF:
            return "PDF解析库未安装"
        
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            if doc.page_count > 0:
                page = doc[0]
                return page.get_text()
            return "空PDF文件"
        except Exception:
            return "PDF解析失败"
    
    @staticmethod
    def extract_docx_content(content):
        """提取DOCX内容"""
        if not HAS_DOCX:
            return "DOCX解析库未安装"
        
        try:
            doc = DocxDocument(io.BytesIO(content))
            return '\n'.join([para.text for para in doc.paragraphs])
        except Exception:
            return "DOCX解析失败"
    
    @staticmethod
    def extract_xlsx_content(content):
        """提取XLSX内容"""
        if not HAS_OPENPYXL:
            return None
        
        try:
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            ws = wb.active
            data = []
            for row in ws.iter_rows(max_row=10, max_col=10):
                row_data = []
                for cell in row:
                    val = cell.value
                    row_data.append(str(val) if val else "")
                data.append(row_data)
            return data
        except Exception:
            return None
    
    @staticmethod
    def extract_xls_content(content):
        """提取XLS内容（旧版Excel格式）"""
        if not HAS_XLRD:
            return None
        
        try:
            # xlrd 2.0+不再支持.xls格式，需要使用xlrd<2.0
            try:
                wb = xlrd.open_workbook(file_contents=content)
            except xlrd.XLRDError as e:
                # 尝试其他方式解析
                if "Excel xlsx file" in str(e):
                    # 这实际上是xlsx文件，尝试用openpyxl
                    try:
                        import openpyxl
                        from io import BytesIO
                        wb = openpyxl.load_workbook(BytesIO(content), read_only=True)
                        ws = wb.active
                        data = []
                        for row in ws.iter_rows(max_row=10, max_col=10):
                            row_data = []
                            for cell in row:
                                val = cell.value
                                row_data.append(str(val) if val else "")
                            data.append(row_data)
                        return data
                    except:
                        return None
                return None
            
            ws = wb.sheet_by_index(0)
            data = []
            for i in range(min(10, ws.nrows)):
                row_data = []
                for j in range(min(10, ws.ncols)):
                    val = ws.cell_value(i, j)
                    row_data.append(str(val) if val else "")
                data.append(row_data)
            return data
        except Exception as e:
            # 添加调试信息
            print(f"XLS解析错误: {str(e)}")
            return None
    
    @staticmethod
    def extract_pptx_content(content):
        """提取PPTX内容"""
        if not HAS_PPTX:
            return "PPTX解析库未安装"
        
        try:
            prs = Presentation(io.BytesIO(content))
            if prs.slides:
                text_content = ""
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, 'text'):
                        text_content += shape.text + "\n"
                return text_content.strip()
            return "空PPT文件"
        except Exception:
            return "PPTX解析失败"
    
    @staticmethod
    def extract_doc_content(content):
        """提取DOC内容（旧版Word格式）"""
        try:
            # 尝试使用textract或其他库
            # 如果没有专门的DOC解析库，返回提示信息
            return "DOC文件需要安装python-docx2txt或textract库才能解析"
        except Exception:
            return "DOC解析失败"
    
    @staticmethod
    def extract_csv_content(content):
        """提取CSV内容"""
        try:
            data = []
            content_str = content.decode('utf-8', errors='replace')
            reader = csv.reader(content_str.splitlines())
            for i, row in enumerate(reader):
                if i >= 20:
                    break
                data.append(row)
            return data
        except Exception:
            return None
    
    @staticmethod
    def create_thumbnail(file_type, content, size=(120, 120)):
        """创建文档缩略图"""
        width, height = size
        img = Image.new('RGB', (width, height), color='#f0f0f0')
        draw = ImageDraw.Draw(img)
        
        font = DocumentProcessor.get_font(10)
        
        # 颜色配置
        colors = {
            'PDF': (220, 20, 60),
            'DOCX': (37, 84, 199),
            'DOC': (37, 84, 199),
            'XLSX': (0, 112, 192),
            'XLS': (0, 112, 192),
            'PPTX': (196, 59, 28),
            'PPT': (196, 59, 28)
        }
        color = colors.get(file_type, (100, 100, 100))
        
        # 绘制边框
        draw.rectangle([2, 2, width-2, height-2], outline=color, width=2)
        
        # 绘制文件类型标签
        draw.rectangle([2, 2, width-2, 20], fill=color)
        draw.text((width//2, 7), file_type, fill='white', font=font, anchor='mm')
        
        # 提取并绘制内容预览
        if file_type == 'PDF':
            text = DocumentProcessor.extract_pdf_content(content)[:80]
            DocumentProcessor._draw_text_preview(draw, text, width, height, 28)
        
        elif file_type == 'DOCX':
            text = DocumentProcessor.extract_docx_content(content)[:80]
            DocumentProcessor._draw_text_preview(draw, text, width, height, 28)
        
        elif file_type == 'DOC':
            draw.text((width//2, height//2), "DOC文件", fill=(100, 100, 100), font=font, anchor='mm')
        
        elif file_type == 'XLSX':
            data = DocumentProcessor.extract_xlsx_content(content)
            if data:
                DocumentProcessor._draw_table_preview(draw, data, width, height, 25, font)
            else:
                draw.text((width//2, height//2), "XLSX文件", fill=(100, 100, 100), font=font, anchor='mm')
        
        elif file_type == 'XLS':
            data = DocumentProcessor.extract_xls_content(content)
            if data:
                DocumentProcessor._draw_table_preview(draw, data, width, height, 25, font)
            else:
                draw.text((width//2, height//2), "XLS文件", fill=(100, 100, 100), font=font, anchor='mm')
        
        elif file_type == 'PPTX':
            text = DocumentProcessor.extract_pptx_content(content)[:80]
            DocumentProcessor._draw_text_preview(draw, text, width, height, 28)
        
        return img
    
    @staticmethod
    def _draw_text_preview(draw, text, width, height, start_y):
        """绘制文本预览"""
        if not text:
            draw.text((width//2, height//2), "无内容", fill=(150, 150, 150), 
                      font=DocumentProcessor.get_font(10), anchor='mm')
            return
        
        lines = text.replace('\r', '').split('\n')
        lines = [line.strip() for line in lines if line.strip()][:5]
        
        font = DocumentProcessor.get_font(8)
        y = start_y
        for line in lines:
            if y >= height - 5:
                break
            if len(line) > 12:
                line = line[:12] + "..."
            draw.text((width//2, y), line, fill=(30, 30, 30), font=font, anchor='mm')
            y += 16
    
    @staticmethod
    def _draw_table_preview(draw, data, width, height, start_y, font):
        """绘制表格预览"""
        start_x, cell_w, cell_h = 5, 28, 15
        max_rows, max_cols = 5, 4
        
        # 绘制边框
        draw.rectangle([start_x, start_y, 
                        start_x + cell_w * max_cols + 1, 
                        start_y + cell_h * max_rows + 1], 
                       outline=(180, 180, 180), width=1)
        
        # 绘制单元格
        for i in range(min(max_rows, len(data))):
            for j in range(min(max_cols, len(data[i]))):
                cell_x = start_x + j * cell_w
                cell_y = start_y + i * cell_h
                
                draw.rectangle([cell_x, cell_y, cell_x + cell_w, cell_y + cell_h], 
                               outline=(200, 200, 200), width=1)
                
                text = str(data[i][j])[:6]
                if text:
                    draw.text((cell_x + cell_w//2, cell_y + cell_h//2), 
                              text, fill=(30, 30, 30), font=DocumentProcessor.get_font(7), anchor='mm')
    
    @staticmethod
    def create_large_preview(file_type, content, width=800, height=600):
        """创建大尺寸预览图"""
        preview_width = width - 20
        preview_height = height - 60
        
        img = Image.new('RGB', (preview_width, preview_height), color='#ffffff')
        draw = ImageDraw.Draw(img)
        
        font = DocumentProcessor.get_font(12)
        
        # 颜色配置
        colors = {
            'PDF': (220, 20, 60),
            'DOCX': (37, 84, 199),
            'DOC': (37, 84, 199),
            'XLSX': (0, 112, 192),
            'XLS': (0, 112, 192),
            'PPTX': (196, 59, 28),
            'CSV': (0, 150, 0)
        }
        color = colors.get(file_type, (100, 100, 100))
        
        # 绘制标题栏
        title_height = 35
        draw.rectangle([0, 0, preview_width, title_height], fill=color)
        draw.text((preview_width//2, title_height//2), file_type, fill='white', 
                  font=DocumentProcessor.get_font(12), anchor='mm')
        
        if file_type == 'PDF':
            text = DocumentProcessor.extract_pdf_content(content)
            DocumentProcessor._draw_large_text_preview(draw, text, preview_width, preview_height, title_height)
        
        elif file_type == 'DOCX':
            text = DocumentProcessor.extract_docx_content(content)
            DocumentProcessor._draw_large_text_preview(draw, text, preview_width, preview_height, title_height)
        
        elif file_type == 'DOC':
            text = DocumentProcessor.extract_doc_content(content)
            DocumentProcessor._draw_large_text_preview(draw, text, preview_width, preview_height, title_height)
        
        elif file_type == 'XLSX':
            data = DocumentProcessor.extract_xlsx_content(content)
            if data:
                DocumentProcessor._draw_large_table_preview(draw, data, preview_width, preview_height, title_height)
            else:
                draw.text((preview_width//2, preview_height//2), "XLSX解析失败", 
                          fill=(150, 150, 150), font=DocumentProcessor.get_font(14), anchor='mm')
        
        elif file_type == 'XLS':
            data = DocumentProcessor.extract_xls_content(content)
            if data:
                DocumentProcessor._draw_large_table_preview(draw, data, preview_width, preview_height, title_height)
            else:
                draw.text((preview_width//2, preview_height//2), "XLS解析失败", 
                          fill=(150, 150, 150), font=DocumentProcessor.get_font(14), anchor='mm')
        
        elif file_type == 'CSV':
            data = DocumentProcessor.extract_csv_content(content)
            if data:
                DocumentProcessor._draw_large_table_preview(draw, data, preview_width, preview_height, title_height)
            else:
                draw.text((preview_width//2, preview_height//2), "CSV解析失败", 
                          fill=(150, 150, 150), font=DocumentProcessor.get_font(14), anchor='mm')
        
        elif file_type == 'PPTX':
            text = DocumentProcessor.extract_pptx_content(content)
            DocumentProcessor._draw_large_text_preview(draw, text, preview_width, preview_height, title_height)
        
        return img
    
    @staticmethod
    def _draw_large_text_preview(draw, text_content, width, height, title_height):
        """绘制大尺寸文本预览"""
        if not text_content:
            draw.text((width//2, height//2), "无内容", fill=(150, 150, 150), 
                      font=DocumentProcessor.get_font(14), anchor='mm')
            return
        
        lines = text_content.replace('\r', '').split('\n')
        lines = [line.strip() for line in lines if line.strip()][:25]
        
        font = DocumentProcessor.get_font(11)
        y = title_height + 10
        line_height = 22
        
        for line in lines:
            if y >= height - 20:
                break
            if len(line) > 70:
                line = line[:70] + "..."
            draw.text((width//2, y), line, fill=(30, 30, 30), font=font, anchor='mm')
            y += line_height
    
    @staticmethod
    def _draw_large_table_preview(draw, data, width, height, title_height):
        """绘制大尺寸表格预览"""
        start_x, start_y = 15, title_height + 10
        available_width = width - start_x - 20
        available_height = height - start_y - 20
        
        max_cols = max(4, min(10, int(available_width / 80)))
        max_rows = max(6, min(20, int(available_height / 24)))
        
        cell_w = int(available_width / max_cols)
        cell_h = int(available_height / max_rows)
        
        # 绘制边框
        table_w = cell_w * max_cols + 2
        table_h = cell_h * max_rows + 2
        draw.rectangle([start_x, start_y, start_x + table_w, start_y + table_h], 
                       outline=(150, 150, 150), width=2)
        
        font_size = max(7, min(12, int(cell_h * 0.4)))
        font = DocumentProcessor.get_font(font_size)
        
        # 绘制单元格
        for i in range(min(max_rows, len(data))):
            for j in range(min(max_cols, len(data[i]))):
                cell_x = start_x + j * cell_w
                cell_y = start_y + i * cell_h
                
                draw.rectangle([cell_x, cell_y, cell_x + cell_w, cell_y + cell_h], 
                               outline=(200, 200, 200), width=1)
                
                text = str(data[i][j])[:int(cell_w/6)]
                if len(text) > int(cell_w/6):
                    text = text[:int(cell_w/6)-3] + "..."
                
                if text:
                    draw.text((cell_x + cell_w//2, cell_y + cell_h//2), 
                              text, fill=(30, 30, 30), font=font, anchor='mm')
